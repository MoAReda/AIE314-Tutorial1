# Document Processing and Retrieval-Augmented Generation (RAG)

# Install Required Libraries
#Run this line in cmd 
#pip install python-docx docx2txt python-pptx PyPDF2 pandas requests beautifulsoup4 transformers faiss-cpu tesseract pytesseract ebooklib
#pip install sentence-transformers pytesseract

# Import Libraries
import docx2txt
import pptx
import PyPDF2
import pandas as pd
import requests
from bs4 import BeautifulSoup
import re
from transformers import pipeline
import os
import mimetypes
import faiss
from sentence_transformers import SentenceTransformer
import pytesseract
from PIL import Image
from ebooklib import epub
import json

# Document Extraction Functions

def extract_text_from_docx(file_path):
    return docx2txt.process(file_path)

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame])

def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        return '\n'.join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path)
    return '\n'.join(df.astype(str).values.flatten())

def extract_text_from_epub(file_path):
    book = epub.read_epub(file_path)
    return '\n'.join([item.get_body_content().decode('utf-8') for item in book.get_items() if item.get_type() == 9])

def extract_text_from_image(file_path):
    return pytesseract.image_to_string(Image.open(file_path))

def read_html_with_requests(url):
    response = requests.get(url)
    response.raise_for_status()
    return BeautifulSoup(response.content, 'html.parser').get_text()

# File Type Identification
def get_file_type(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if 'msword' in mime_type:
            return 'docx'
        elif 'presentation' in mime_type:
            return 'pptx'
        elif 'pdf' in mime_type:
            return 'pdf'
        elif 'excel' in mime_type:
            return 'excel'
        elif 'epub' in mime_type:
            return 'epub'
        elif mime_type.startswith('image/'):
            return 'image'
    return 'unknown'

# Convert to JSON and Save to Output Folder
def convert_to_json(file_type, content, path):
    json_data = {"file_type": file_type, "content": content, "meta_data": {"author": "unknown", "date": "unknown", "path": path}}
    
    # Ensure output directory exists
    os.makedirs('output', exist_ok=True)
    
    # Save JSON file
    output_path = os.path.join('output', f'{file_type}_output.json')
    with open(output_path, 'w', encoding='utf-8') as json_file:
        json.dump(json_data, json_file, ensure_ascii=False, indent=4)
        
    print(f"JSON file saved at: {output_path}")
    return json_data

# **RAG Component**

# Initialize Sentence Transformer and FAISS
embedder = SentenceTransformer('all-MiniLM-L6-v2')
index = None

def build_vector_store(text_chunks):
    global index
    embeddings = embedder.encode(text_chunks)
    dimension = embeddings.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings)
    return text_chunks

def retrieve_relevant_chunks(query, text_chunks, top_k=3):
    query_embedding = embedder.encode([query])
    _, indices = index.search(query_embedding, top_k)
    return [text_chunks[i] for i in indices[0]]

# Question Answering with Transformers
qa_pipeline = pipeline("question-answering")

def answer_question_rag(question, context_chunks):
    context = ' '.join(context_chunks)
    return qa_pipeline(question=question, context=context)['answer']

# Example Usage
if __name__ == "__main__":
    file_path = "ML.png"
    file_type = get_file_type(file_path)

    if file_type == 'docx':
        content = extract_text_from_docx(file_path)
    elif file_type == 'pptx':
        content = extract_text_from_pptx(file_path)
    elif file_type == 'pdf':
        content = extract_text_from_pdf(file_path)
    elif file_type == 'excel':
        content = extract_text_from_excel(file_path)
    elif file_type == 'epub':
        content = extract_text_from_epub(file_path)
    elif file_type == 'image':
        content = extract_text_from_image(file_path)
    else:
        content = ""

    if content:
        text_chunks = content.split('\n')
        text_chunks = build_vector_store(text_chunks)

        json_output = convert_to_json(file_type, content, file_path)
        print("--- JSON Output ---")
        print(json_output)

        question = "What is the topic?"
        print("--- Question ---")
        print(question)

        relevant_chunks = retrieve_relevant_chunks(question, text_chunks)
        print("--- Relevant Chunks ---")
        for chunk in relevant_chunks:
            print(chunk)

        answer = answer_question_rag(question, relevant_chunks)
        print("--- Answer to Question ---")
        print(answer)
    else:
        print("Unsupported file type or empty content.")
